import json

from pptx import Presentation

from flash_cards import FlashCards
from question_generator import QuestionGenerator
from question_model import Question

# Load the syllabus PowerPoint file
try:
    syllabus = Presentation("syllabus/example_syllabus.pptx")
except Exception as e:
    print("Error loading presentation:", e)
    exit()

# Extract text from all slides in the PowerPoint
syllabus_text = ""
for slide in syllabus.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):  # Only extract if the shape has text content
            syllabus_text += shape.text + "\n"

# Create an instance of the question generator
question_generator = QuestionGenerator()

# Generate 10 questions from the syllabus text
questions = question_generator.get_questions(syllabus_text, 10)

# Convert the generated JSON string into Python objects
try:
    question_list = json.loads(questions)
    question_list = [Question(question) for question in question_list]
except json.JSONDecodeError as e:
    print("Error parsing generated questions:", e)
    exit()

# Start the flashcard GUI with the list of questions
FlashCards(question_list)
